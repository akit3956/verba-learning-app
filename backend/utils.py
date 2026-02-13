from pptx import Presentation
import openpyxl
from io import BytesIO

def extract_text_from_pptx(file_content: bytes) -> str:
    prs = Presentation(BytesIO(file_content))
    text = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
                
    return "\n".join(text)

def extract_text_from_xlsx(file_content: bytes) -> str:
    workbook = openpyxl.load_workbook(BytesIO(file_content), data_only=True)
    text = []
    
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            # Filter None and convert to string
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                text.append(" ".join(row_text))
                
    return "\n".join(text)
