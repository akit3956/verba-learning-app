import os
import glob
import json
import pdfplumber
from docx import Document
from pptx import Presentation
from openai import OpenAI
import psycopg2
from psycopg2.extras import Json
from dotenv import load_dotenv

load_dotenv()

# Configuration
NOTES_DIR = os.path.join(os.path.dirname(__file__), "teacher_notes")
PROJECT_ROOT = os.path.dirname(os.path.dirname(__file__)) # One level up from backend
DB_URL = os.getenv("DATABASE_URL")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

client = OpenAI(api_key=OPENAI_API_KEY)

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    try:
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        elif ext == ".pdf":
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += (page.extract_text() or "") + "\n"
        elif ext == ".docx":
            doc = Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting {file_path}: {e}")
    return text

def extract_excel_text(file_path):
    import openpyxl
    text = ""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet in wb.worksheets:
            text += f"--- Sheet: {sheet.title} ---\n"
            for row in sheet.iter_rows(values_only=True):
                # Filter out None and join values with space
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    text += row_text + "\n"
    except Exception as e:
        print(f"Error extracting Excel {file_path}: {e}")
    return text

def extract_text_ext(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".xlsx":
        return extract_excel_text(file_path)
    return extract_text(file_path) # Fallback to original

def chunk_text(text, chunk_size=800, overlap=100):
    chunks = []
    # Simple character-based chunking for now
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += (chunk_size - overlap)
    return chunks

def get_embedding(text):
    response = client.embeddings.create(
        input=text,
        model="text-embedding-3-small"
    )
    return response.data[0].embedding

def ingest():
    if not DB_URL:
        print("DATABASE_URL not found.")
        return

    conn = psycopg2.connect(DB_URL)
    cur = conn.cursor()

    files = glob.glob(os.path.join(NOTES_DIR, "*.*"))
    root_file = os.path.join(PROJECT_ROOT, "Verba N3 模擬試験.docx")
    if os.path.exists(root_file):
        files.append(root_file)
        
    print(f"Found {len(files)} files to process.")

    # Get already ingested filenames from DB
    cur.execute("SELECT DISTINCT metadata->>'source' FROM teacher_embeddings")
    ingested_sources = {row[0] for row in cur.fetchall()}
    print(f"Already ingested {len(ingested_sources)} sources.")

    for file_path in files:
        file_name = os.path.basename(file_path)
        
        if file_name in ingested_sources:
            print(f"Skipping {file_name} (already ingested).")
            continue

        print(f"Processing {file_name}...")
        full_text = extract_text_ext(file_path)
        if not full_text:
            continue
        
        chunks = chunk_text(full_text)
        print(f"  Split into {len(chunks)} chunks.")

        for i, chunk in enumerate(chunks):
            embedding = get_embedding(chunk)
            metadata = {
                "source": os.path.basename(file_path),
                "chunk_index": i
            }
            
            cur.execute(
                "INSERT INTO teacher_embeddings (content, metadata, embedding) VALUES (%s, %s, %s)",
                (chunk, Json(metadata), embedding)
            )
        
        conn.commit()
    
    cur.close()
    conn.close()
    print("Ingestion complete.")

if __name__ == "__main__":
    ingest()
